import tkinter as tk
from tkinter import filedialog
import xlwings as xw
import os
from pptx import Presentation
from pptx.util import Inches
from docx import Document

def load_excel():
    global excel_file_path
    excel_file_path = filedialog.askopenfilename(title="选择Excel文件", filetypes=[("Excel files", "*.xlsx;*.xlsm;*.xls")])
    if excel_file_path:
        load_data()

def load_data():
    # 打开Excel文件
    wb = xw.Book(excel_file_path)

    # 创建PPT和Word文件
    ppt = Presentation()
    doc = Document()

    # 遍历每个表单
    for sheet in wb.sheets:
        charts = sheet.api.ChartObjects()  # 获取所有图表
        for chart in charts:
            chart_obj = chart.Chart
            # 导出图表为图片
            chart_path = os.path.join(os.path.dirname(excel_file_path), f"{sheet.name}_{chart.Index}.png")
            chart_obj.Export(chart_path)

            # 添加图片到PPT
            slide_layout = ppt.slide_layouts[5]  # 空白布局
            slide = ppt.slides.add_slide(slide_layout)
            left = Inches(1)
            top = Inches(2)
            pic = slide.shapes.add_picture(chart_path, left, top, width=Inches(7.5), height=Inches(4.5))

            # 添加图片到Word
            doc.add_picture(chart_path, width=Inches(5))

            # 删除临时图片文件
            os.remove(chart_path)

    # 保存PPT和Word文件
    output_dir = os.path.dirname(excel_file_path)
    ppt_output_path = os.path.join(output_dir, "output_presentation.pptx")
    word_output_path = os.path.join(output_dir, "output_document.docx")

    ppt.save(ppt_output_path)
    doc.save(word_output_path)

    ppt_info_label.config(text=f"PPT文件已保存到EXCEL所在文件夹")
    word_info_label.config(text=f"Word文件已保存到EXCEL所在文件夹")



# 创建GUI窗口
root = tk.Tk()
root.title("作者：审协江苏中心光电部尹文杰")
root.geometry("400x250")

# 添加标题栏
title_label = tk.Label(root, text="Excel图像批量至WORD/PPT工具", font=("Arial", 18))
title_label.pack(pady=10)

# 加载Excel文件按钮
load_button = tk.Button(root, text="加载Excel文件", command=load_excel)
load_button.pack(pady=20)

# 显示保存信息的标签
ppt_info_label = tk.Label(root, text="", font=("Arial", 10))
ppt_info_label.pack(pady=5)

word_info_label = tk.Label(root, text="", font=("Arial", 10))
word_info_label.pack(pady=5)

root.mainloop()
